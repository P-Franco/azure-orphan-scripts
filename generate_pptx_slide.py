#!/usr/bin/env python3
"""
generate_pptx_slide.py

Generates a single CIR-ready PowerPoint slide summarising orphaned Azure
resources.  Reads from a JSON report produced by orphan_report.py.

Usage:
  python3 orphan_report.py --format json --output scan.json
  python3 generate_pptx_slide.py --input scan.json
  python3 generate_pptx_slide.py --input scan.json --output monthly-cir.pptx --client "Contoso Corp"
"""

import argparse
import json
import sys
from collections import Counter
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Colour palette ───────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)       # dark navy
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)   # bright blue
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # emerald
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # amber
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)     # red
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)     # slate-400
CARD_BG = RGBColor(0x1E, 0x33, 0x5C)        # slightly lighter navy
TABLE_HEADER_BG = RGBColor(0x33, 0x4E, 0x7E)
TABLE_ROW_ALT = RGBColor(0x24, 0x3B, 0x67)


def _add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle shape as a card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Subtle rounding
    shape.adjustments[0] = 0.05
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=12,
                  bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Add a text box with the specified formatting."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def _add_kpi_card(slide, left, top, label, value, color=ACCENT_BLUE):
    """Add a KPI card (label + big number)."""
    card_w = Inches(2.7)
    card_h = Inches(1.35)
    _add_rounded_rect(slide, left, top, card_w, card_h, CARD_BG)
    # Value (big number)
    _add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                  card_w - Inches(0.4), Inches(0.7),
                  str(value), font_size=32, bold=True, color=color,
                  alignment=PP_ALIGN.CENTER)
    # Label
    _add_text_box(slide, left + Inches(0.2), top + Inches(0.85),
                  card_w - Inches(0.4), Inches(0.4),
                  label, font_size=11, bold=False, color=LIGHT_GRAY,
                  alignment=PP_ALIGN.CENTER)


def _build_slide(prs, data, client_name):
    """Build the single CIR summary slide."""
    # Use a blank layout
    blank_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Dark background ──────────────────────────────────────────────────
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # ── Parse data ───────────────────────────────────────────────────────
    resources = data.get("resources", [])
    total = data.get("totalResources", len(resources))
    total_cost = data.get("estimatedMonthlyCost",
                          sum(r.get("estimatedMonthlyCost", 0) for r in resources))
    scan_date = data.get("generatedAt", datetime.now(timezone.utc).isoformat())

    # Parse date for display
    try:
        dt = datetime.fromisoformat(scan_date.replace("Z", "+00:00"))
        date_str = dt.strftime("%B %Y")
    except Exception:
        date_str = scan_date[:10]

    prod_count = sum(1 for r in resources if r.get("environment") == "PRODUCTION")
    nonprod_count = total - prod_count

    # Top resource types
    type_counts = Counter(r.get("category", "Unknown") for r in resources)
    top_types = type_counts.most_common(6)

    # Subscriptions scanned
    subs = set(r.get("subscription", "") for r in resources)
    sub_count = len(subs) if subs - {""} else 0

    # ── Title bar ────────────────────────────────────────────────────────
    title_text = "Azure Orphaned Resources — Monthly Summary"
    _add_text_box(slide, Inches(0.5), Inches(0.25), Inches(7), Inches(0.45),
                  title_text, font_size=20, bold=True, color=WHITE)

    subtitle = f"{client_name}  •  {date_str}"
    _add_text_box(slide, Inches(0.5), Inches(0.65), Inches(7), Inches(0.3),
                  subtitle, font_size=11, bold=False, color=LIGHT_GRAY)

    # ── KPI Cards (row of 3) ─────────────────────────────────────────────
    kpi_top = Inches(1.15)
    kpi_gap = Inches(0.25)
    kpi_w = Inches(2.7)

    _add_kpi_card(slide, Inches(0.5), kpi_top,
                  "Total Orphaned Resources", f"{total:,}", ACCENT_BLUE)
    _add_kpi_card(slide, Inches(0.5) + kpi_w + kpi_gap, kpi_top,
                  "Est. Monthly Waste", f"${total_cost:,.2f}", ACCENT_RED)
    _add_kpi_card(slide, Inches(0.5) + 2 * (kpi_w + kpi_gap), kpi_top,
                  "Subscriptions Scanned", f"{sub_count}", ACCENT_GREEN)

    # ── Left section: Env breakdown ──────────────────────────────────────
    section_top = Inches(2.75)

    # Environment card
    env_card_w = Inches(3.8)
    env_card_h = Inches(1.3)
    _add_rounded_rect(slide, Inches(0.5), section_top, env_card_w, env_card_h, CARD_BG)

    _add_text_box(slide, Inches(0.7), section_top + Inches(0.1),
                  Inches(3.4), Inches(0.3),
                  "ENVIRONMENT SPLIT", font_size=10, bold=True, color=LIGHT_GRAY)

    # Production bar
    bar_left = Inches(0.7)
    bar_top = section_top + Inches(0.5)
    bar_full_w = Inches(3.4)
    bar_h = Inches(0.28)

    prod_pct = (prod_count / total * 100) if total > 0 else 0
    nonprod_pct = 100 - prod_pct

    # Prod bar
    if prod_count > 0:
        prod_w = int(bar_full_w * prod_pct / 100)
        prod_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_left, bar_top, prod_w, bar_h
        )
        prod_shape.fill.solid()
        prod_shape.fill.fore_color.rgb = ACCENT_GREEN
        prod_shape.line.fill.background()
        prod_shape.adjustments[0] = 0.3

    # Non-prod bar
    if nonprod_count > 0:
        nonprod_w = int(bar_full_w * nonprod_pct / 100)
        nonprod_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_left + int(bar_full_w * prod_pct / 100), bar_top,
            nonprod_w, bar_h
        )
        nonprod_shape.fill.solid()
        nonprod_shape.fill.fore_color.rgb = ACCENT_AMBER
        nonprod_shape.line.fill.background()
        nonprod_shape.adjustments[0] = 0.3

    # Legend
    _add_text_box(slide, Inches(0.7), bar_top + Inches(0.35),
                  Inches(1.7), Inches(0.25),
                  f"● Production: {prod_count} ({prod_pct:.0f}%)",
                  font_size=9, color=ACCENT_GREEN)
    _add_text_box(slide, Inches(2.4), bar_top + Inches(0.35),
                  Inches(1.7), Inches(0.25),
                  f"● Non-Prod: {nonprod_count} ({nonprod_pct:.0f}%)",
                  font_size=9, color=ACCENT_AMBER)

    # ── Right section: Top resource types table ──────────────────────────
    tbl_left = Inches(4.6)
    tbl_top = section_top
    tbl_w = Inches(4.9)
    num_rows = min(len(top_types), 6) + 1  # +1 for header
    row_h = Inches(0.28)
    tbl_h = row_h * num_rows

    _add_rounded_rect(slide, tbl_left - Inches(0.1), tbl_top,
                      tbl_w + Inches(0.2), Inches(1.3), CARD_BG)

    _add_text_box(slide, tbl_left, tbl_top + Inches(0.1),
                  tbl_w, Inches(0.3),
                  "TOP RESOURCE TYPES", font_size=10, bold=True, color=LIGHT_GRAY)

    # Table rows
    y = tbl_top + Inches(0.42)
    for i, (cat_name, count) in enumerate(top_types):
        # Truncate long names
        display_name = cat_name[:38] + "…" if len(cat_name) > 38 else cat_name
        _add_text_box(slide, tbl_left, y, Inches(3.8), Inches(0.22),
                      display_name, font_size=9, color=WHITE)
        _add_text_box(slide, tbl_left + Inches(3.8), y, Inches(0.9), Inches(0.22),
                      str(count), font_size=9, bold=True, color=ACCENT_BLUE,
                      alignment=PP_ALIGN.RIGHT)
        y += Inches(0.22)

    # ── Bottom: Recommendation strip ─────────────────────────────────────
    rec_top = Inches(4.3)
    _add_rounded_rect(slide, Inches(0.5), rec_top, Inches(9.0), Inches(0.55), CARD_BG)

    # Build recommendation text
    if total_cost > 0:
        rec_text = (f"💡 Recommendation: Review and clean up {total:,} orphaned resources "
                    f"to save an estimated ${total_cost:,.2f}/month. "
                    f"Run: python3 orphan_cleanup.py --production-only --dry-run")
    else:
        rec_text = (f"💡 {total:,} orphaned resources found with no direct cost impact. "
                    f"Clean up recommended for hygiene and security posture.")

    _add_text_box(slide, Inches(0.7), rec_top + Inches(0.12),
                  Inches(8.6), Inches(0.35),
                  rec_text, font_size=9, color=WHITE)

    # ── Footer ───────────────────────────────────────────────────────────
    footer = f"Generated by azure-orphan-scripts  •  {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
    _add_text_box(slide, Inches(0.5), Inches(5.05), Inches(9.0), Inches(0.25),
                  footer, font_size=7, color=LIGHT_GRAY,
                  alignment=PP_ALIGN.RIGHT)


def main():
    parser = argparse.ArgumentParser(
        description="Generate a CIR-ready PowerPoint slide from orphan report JSON"
    )
    parser.add_argument("--input", "-i", required=True,
                        help="Path to JSON report from orphan_report.py")
    parser.add_argument("--output", "-o", default=None,
                        help="Output .pptx file (default: orphan-cir-slide-YYYYMM.pptx)")
    parser.add_argument("--client", "-c", default="Azure Tenant",
                        help="Client name for the slide header")
    args = parser.parse_args()

    # Load JSON
    try:
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    except FileNotFoundError:
        print(f"Error: File not found: {args.input}")
        return 1
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON: {e}")
        return 1

    # Output path
    if args.output:
        out_path = args.output
    else:
        month_str = datetime.now(timezone.utc).strftime("%Y%m")
        out_path = f"orphan-cir-slide-{month_str}.pptx"

    # Build presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    _build_slide(prs, data, args.client)

    prs.save(out_path)
    print(f"✓ CIR slide saved to: {out_path}")
    print(f"  Client: {args.client}")
    print(f"  Resources: {data.get('totalResources', len(data.get('resources', [])))}")
    print(f"  Est. waste: ${data.get('estimatedMonthlyCost', 0):,.2f}/month")
    return 0


if __name__ == "__main__":
    sys.exit(main() or 0)
